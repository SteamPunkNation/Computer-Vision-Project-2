from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

from pptx.dml.color import RGBColor

def add_slide(prs, title, content, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(34, 34, 34)
    
    # Style Title
    title_shape = slide.shapes.title
    title_shape.text = title
    if title_shape.text_frame.paragraphs:
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(77, 184, 255) # Light blue
        title_shape.text_frame.paragraphs[0].font.name = "Arial"
    
    # Style Content
    tf = slide.placeholders[1].text_frame
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.color.rgb = RGBColor(240, 240, 240)
        paragraph.font.name = "Arial"
    
    # Presenter Notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes

# Slide 1
add_slide(prs, 
          "Semantic-Aware Monocular Visual Odometry", 
          "Running on Local Hardware", 
          "Presenter Instructions:\n- Have your terminal open in the background, ready to run the demo.\n- Stand confidently and make eye contact with the audience.\n\nVerbatim Script:\n\"Hello everyone. Today I am presenting my project on Semantic-Aware Monocular Visual Odometry. The goal of this project was to build a real-time localization system that runs entirely on local laptop hardware, using only a single standard webcam, without the need for expensive tracking arrays or depth sensors.\"")

# Slide 2
add_slide(prs,
          "The Problem",
          "• Scale Ambiguity: A single camera cannot easily determine physical scale.\n• Dynamic Obstacles: Moving objects (people, cars, hands) ruin tracking geometry.\n• Cumulative Drift: Small mathematical errors accumulate over time, destroying the trajectory.",
          "Presenter Instructions:\n- Gesture to emphasize the difficulty of these three computer vision problems.\n\nVerbatim Script:\n\"When building a monocular visual odometry system, we face three massive hurdles. First, a single camera has no concept of depth or scale. Second, moving objects in the frame—like people or hands—trick the math into thinking the camera is moving when it isn't. And third, even perfect algorithms suffer from cumulative drift, where microscopic frame-to-frame errors eventually ruin the entire trajectory plot.\"")

# Slide 3
add_slide(prs,
          "Our Solution",
          "• Semantic Masking: YOLOv8 Nano & MediaPipe Hands mask dynamic objects.\n• Dynamic Scale Estimation: Triangulating ground-plane features using a camera height heuristic.\n• Windowed Bundle Adjustment: Scipy-based backend optimization over a sliding window of frames.",
          "Presenter Instructions:\n- Speak clearly about the specific technologies used.\n\nVerbatim Script:\n\"To solve these issues, I built a modular Python pipeline. First, I integrated a lightweight YOLOv8 segmentation model and MediaPipe to actively mask out dynamic foreground obstacles before tracking begins. To solve scale, the system triangulates ground-plane features and anchors them using a known camera height. Finally, to eliminate drift, I implemented a local Windowed Bundle Adjustment backend using SciPy to continuously optimize our trajectory over a sliding window.\"")

# Slide 4
add_slide(prs,
          "Live Demo",
          "(Switch to terminal and run the application)\n\npython main.py --camera-height 1.2 --window-size 5",
          "Presenter Instructions:\n- Alt/Cmd-Tab to your terminal.\n- Run: python main.py --camera-height 1.2 --window-size 5\n- Move your laptop or webcam deliberately left, right, forward, and backward.\n- Point out the masked objects (red tint) and the trajectory plot updating in real-time.\n\nVerbatim Script:\n\"I'd now like to show you a live demonstration. As you can see, the application is capturing my webcam feed. The red tinted areas are the semantic masks actively ignoring my hands and body. Notice the green tracking points on the static background. As I move the camera left and right, you can see the trajectory plot mapping my physical movement in real-time, scaled to meters, while maintaining a smooth frame rate.\"")

# Slide 5
add_slide(prs,
          "Results & Conclusion",
          "• Performance: Achieved real-time FPS on local hardware.\n• Accuracy: Masking significantly reduces tracking outliers.\n• Future Work: Full backend loop-closure and IMU integration.",
          "Presenter Instructions:\n- Close the demo ('q') and return to the slides.\n\nVerbatim Script:\n\"In conclusion, the integration of deep learning semantic segmentation with classical epipolar geometry proved highly successful. We achieved real-time performance on local hardware, and the Windowed Bundle Adjustment successfully mitigated cumulative drift. Thank you for your time, I am happy to take any questions.\"")

prs.save('presentation.pptx')
